
# -*- coding: utf-8 -*-
"""
Auto Real-Estate PPT Report (with News via RSS)
- Includes Sale/Jeonse/Auction/Population/Supply/Jobs/PIR
- Adds News section fetched from RSS for query "청주시 하이닉스"
- Filters news to items within the last 1 month from today
Notes:
  * Charts: matplotlib (no seaborn), single plot per figure, no explicit colors
  * If RSS fetch fails (e.g., offline), script falls back to mock news items
"""
import os, sys, sqlite3
from datetime import datetime, timedelta
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

try:
    import requests
except Exception:
    requests = None
import xml.etree.ElementTree as ET

BASE_DIR = "/mnt/data"
DB_PATH = os.path.join(BASE_DIR, "fresh_realestate_demo_with_news.db")
PPT_PATH = os.path.join(BASE_DIR, "Auto_RealEstate_Report_With_News.pptx")

def fetch_rss_news(query: str, max_items: int = 8, days: int = 31):
    cutoff = datetime.now() - timedelta(days=days)
    items = []
    q = query.replace(" ", "+")
    url = f"https://news.google.com/rss/search?q={q}&hl=ko&gl=KR&ceid=KR:ko"
    try:
        if requests is None:
            raise RuntimeError("requests not available")
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        root = ET.fromstring(resp.content)
        for item in root.findall(".//item"):
            title = (item.findtext("title") or "").strip()
            link = (item.findtext("link") or "").strip()
            source = ""
            src_el = item.find("{http://search.yahoo.com/mrss/}source")
            if src_el is not None and src_el.text:
                source = src_el.text.strip()
            else:
                if " - " in title:
                    parts = title.rsplit(" - ", 1)
                    title, source = parts[0], parts[1]
            pubdate_txt = item.findtext("pubDate") or ""
            try:
                pubdate = datetime.strptime(pubdate_txt[:25], "%a, %d %b %Y %H:%M:%S")
            except Exception:
                pubdate = datetime.now()
            if pubdate >= cutoff:
                items.append({"title": title, "source": source or "뉴스", "pubdate": pubdate, "link": link})
            if len(items) >= max_items:
                break
    except Exception:
        now = datetime.now()
        items = [{
            "title": f"[모의] 청주시 하이닉스 관련 투자·고용 소식 {i+1}",
            "source": "모의뉴스",
            "pubdate": now - timedelta(days=3*i+1),
            "link": "https://example.com/news"
        } for i in range(max_items)]
    items.sort(key=lambda x: x["pubdate"], reverse=True)
    return items[:max_items]

# --- Create demo DB ---
conn = sqlite3.connect(DB_PATH)
cur = conn.cursor()
for t in ["transactions", "auctions", "population_trend", "supply", "jobs", "income"]:
    cur.execute(f"DROP TABLE IF EXISTS {t}")

cur.execute("""CREATE TABLE transactions(
  id INTEGER PRIMARY KEY AUTOINCREMENT,
  deal_date TEXT, address TEXT, apt_name TEXT, area_m2 REAL,
  sale_price_million INTEGER, jeonse_million INTEGER)""")

cur.execute("""CREATE TABLE auctions(
  id INTEGER PRIMARY KEY AUTOINCREMENT,
  auction_date TEXT, address TEXT, apt_name TEXT, area_m2 REAL,
  appraised_million INTEGER, hammer_million INTEGER)""")

cur.execute("""CREATE TABLE population_trend(ym TEXT, dong TEXT, population INTEGER)""")
cur.execute("""CREATE TABLE supply(year INTEGER, area TEXT, completions INTEGER)""")
cur.execute("""CREATE TABLE jobs(ym TEXT, area TEXT, employed INTEGER)""")
cur.execute("""CREATE TABLE income(ym TEXT, area TEXT, annual_income_million REAL)""")

target_address = "충북 청주시 서원구 개신동 1-1"
nearby_flag = "충북 청주시 서원구 개신동"
dates = pd.date_range(end=datetime.today(), periods=12, freq="MS")
rng = np.random.default_rng(20250810)

tx_rows = []
base_sale = 20000
for i, d in enumerate(dates):
    sale1 = base_sale + i*260 + int(rng.integers(-120, 150))
    sale2 = sale1 + 1300
    jr1 = 0.67 + rng.random()*0.11
    jr2 = 0.67 + rng.random()*0.11
    tx_rows += [
        (d.strftime("%Y-%m-%d"), f"{nearby_flag} 101-{i+1}", "개신 힐스", 59.8, int(sale1), int(sale1*jr1)),
        (d.strftime("%Y-%m-%d"), f"{nearby_flag} 202-{i+1}", "개신 더샵", 84.9, int(sale2), int(sale2*jr2)),
    ]
cur.executemany("""INSERT INTO transactions(deal_date,address,apt_name,area_m2,sale_price_million,jeonse_million)
VALUES(?,?,?,?,?,?)""", tx_rows)

auc_rows = []
for i in range(6):
    dt = pd.to_datetime(dates[-1]) - pd.DateOffset(days=10*i)
    app = 26000 + int(rng.integers(-900, 900))
    ham = int(app * (0.82 + rng.random()*0.09))
    auc_rows.append((dt.strftime("%Y-%m-%d"), f"{nearby_flag} 30{i}", "개신 리버뷰", 84.9, int(app), int(ham)))
cur.executemany("""INSERT INTO auctions(auction_date,address,apt_name,area_m2,appraised_million,hammer_million)
VALUES(?,?,?,?,?,?)""", auc_rows)

pop_rows = []
pop_base = 18600
for i in range(24):
    ym = (pd.to_datetime(dates[0]) - pd.DateOffset(months=12) + pd.DateOffset(months=i)).strftime("%Y-%m")
    pop_rows.append((ym, "개신동", int(pop_base + i*18 + int(rng.integers(-50, 50)))))
cur.executemany("INSERT INTO population_trend(ym,dong,population) VALUES(?,?,?)", pop_rows)

supply_rows = []
for y in range(datetime.today().year-2, datetime.today().year+3):
    supply_rows.append((y, "청주시 서원구", int(rng.integers(250, 1300))))
cur.executemany("INSERT INTO supply(year,area,completions) VALUES(?,?,?)", supply_rows)

job_rows = []
for i in range(12):
    ym = (pd.to_datetime(dates[-1]) - pd.DateOffset(months=11-i)).strftime("%Y-%m")
    job_rows.append((ym, "청주시", int(241000 + i*520 + int(rng.integers(-1500, 1500)))))
cur.executemany("INSERT INTO jobs(ym,area,employed) VALUES(?,?,?)", job_rows)

income_rows = []
base_income = 5250
for i in range(12):
    ym = (pd.to_datetime(dates[-1]) - pd.DateOffset(months=11-i)).strftime("%Y-%m")
    val = base_income + i*28 + int(rng.integers(-30, 30))
    income_rows.append((ym, "청주시", float(val)))
cur.executemany("INSERT INTO income(ym,area,annual_income_million) VALUES(?,?,?)", income_rows)

conn.commit()

# --- Query & KPIs ---
tx_df = pd.read_sql_query("""SELECT substr(deal_date,1,7) AS ym, apt_name, area_m2, sale_price_million, jeonse_million
FROM transactions WHERE address LIKE ?""", conn, params=(f"%{nearby_flag}%",))

auc_df = pd.read_sql_query("""SELECT auction_date, apt_name, area_m2, appraised_million, hammer_million,
ROUND(100.0*hammer_million/appraised_million,1) AS rate
FROM auctions WHERE address LIKE ? ORDER BY auction_date ASC""", conn, params=(f"%{nearby_flag}%",))

pop_df = pd.read_sql_query("""SELECT ym, population FROM population_trend WHERE dong = '개신동' ORDER BY ym""", conn)
sup_df = pd.read_sql_query("""SELECT year, completions FROM supply WHERE area = '청주시 서원구' ORDER BY year""", conn)
jobs_df = pd.read_sql_query("""SELECT ym, employed FROM jobs WHERE area = '청주시' ORDER BY ym""", conn)
income_df = pd.read_sql_query("""SELECT ym, annual_income_million FROM income WHERE area = '청주시' ORDER BY ym""", conn)

tx_grp = (tx_df.groupby("ym")[["sale_price_million","jeonse_million"]].mean().reset_index().sort_values("ym"))
tx_grp["jeonse_ratio_pct"] = (tx_grp["jeonse_million"]/tx_grp["sale_price_million"]*100).round(1)

pir_df = pd.merge(tx_grp[["ym","sale_price_million"]], income_df, on="ym", how="inner")
pir_df["PIR"] = (pir_df["sale_price_million"]/pir_df["annual_income_million"]).round(2)

def last_delta(s): return float(s.iloc[-1]-s.iloc[-2]) if len(s)>1 else 0.0
kpi = {
    "월간_평균_매매가(만원)": int(tx_grp["sale_price_million"].iloc[-1]),
    "월간_평균_전세가(만원)": int(tx_grp["jeonse_million"].iloc[-1]),
    "전세가율_최근월(%)": float(tx_grp["jeonse_ratio_pct"].iloc[-1]),
    "매매가_전월대비(만원)": int(last_delta(tx_grp["sale_price_million"])),
    "전세가_전월대비(만원)": int(last_delta(tx_grp["jeonse_million"])),
    "최근_경매_평균_낙찰가율(%)": float(auc_df["rate"].mean()) if not auc_df.empty else None,
    "인구_최근월(명)": int(pop_df["population"].iloc[-1]) if not pop_df.empty else None,
    "연간_공급물량(최근연도,세대)": int(sup_df.set_index("year")["completions"].iloc[-1]) if not sup_df.empty else None,
    "취업자_최근월(명)": int(jobs_df["employed"].iloc[-1]) if not jobs_df.empty else None,
    "PIR_최근월(배)": float(pir_df["PIR"].iloc[-1]) if not pir_df.empty else None,
    "PIR_전월대비(배)": float(last_delta(pir_df["PIR"])) if len(pir_df)>1 else 0.0,
}

# --- Charts ---
IMG_PRICE_BOTH = os.path.join(BASE_DIR, "news_sale_vs_jeonse.png")
IMG_JEONSE_RATIO = os.path.join(BASE_DIR, "news_jeonse_ratio.png")
IMG_PIR = os.path.join(BASE_DIR, "news_pir.png")
IMG_POP = os.path.join(BASE_DIR, "news_population.png")
IMG_SUP = os.path.join(BASE_DIR, "news_supply.png")
IMG_JOBS = os.path.join(BASE_DIR, "news_jobs.png")
IMG_AUC = os.path.join(BASE_DIR, "news_auction.png")

plt.figure(); plt.plot(tx_grp["ym"], tx_grp["sale_price_million"], label="매매가(만원)")
plt.plot(tx_grp["ym"], tx_grp["jeonse_million"], label="전세가(만원)"); plt.legend()
plt.title("월별 평균 매매가 vs 전세가 (만원)"); plt.xticks(rotation=45, ha='right'); plt.tight_layout(); plt.savefig(IMG_PRICE_BOTH, dpi=150); plt.close()

plt.figure(); plt.plot(tx_grp["ym"], tx_grp["jeonse_ratio_pct"])
plt.title("전세가율 추이(%)"); plt.xticks(rotation=45, ha='right'); plt.tight_layout(); plt.savefig(IMG_JEONSE_RATIO, dpi=150); plt.close()

plt.figure(); plt.plot(pir_df["ym"], pir_df["PIR"])
plt.title("PIR(Price-to-Income Ratio) 추이(배)"); plt.xticks(rotation=45, ha='right'); plt.tight_layout(); plt.savefig(IMG_PIR, dpi=150); plt.close()

plt.figure(); plt.plot(pop_df["ym"], pop_df["population"])
plt.title("개신동 인구 추이(명)"); plt.xticks(rotation=45, ha='right'); plt.tight_layout(); plt.savefig(IMG_POP, dpi=150); plt.close()

plt.figure(); plt.bar(sup_df["year"].astype(str), sup_df["completions"])
plt.title("연간 아파트 준공 물량(세대) - 청주시 서원구"); plt.xticks(rotation=0); plt.tight_layout(); plt.savefig(IMG_SUP, dpi=150); plt.close()

plt.figure(); plt.plot(jobs_df["ym"], jobs_df["employed"])
plt.title("청주시 취업자 추이(명)"); plt.xticks(rotation=45, ha='right'); plt.tight_layout(); plt.savefig(IMG_JOBS, dpi=150); plt.close()

plt.figure()
if not auc_df.empty:
    auc_dates = pd.to_datetime(auc_df["auction_date"]); plt.plot(auc_dates, auc_df["rate"])
plt.title("경매 낙찰가율 추이(%)"); plt.xticks(rotation=45, ha='right'); plt.tight_layout(); plt.savefig(IMG_AUC, dpi=150); plt.close()

# --- News (last 1 month) ---
news_items = fetch_rss_news("청주시 하이닉스", max_items=8, days=31)

# --- PPT ---
prs = Presentation()
t = prs.slides.add_slide(prs.slide_layouts[0])
t.shapes.title.text = "자동 부동산 리포트 (News 포함)"
subtitle = t.placeholders[1].text_frame
subtitle.text = "대상지: " + target_address
subtitle.add_paragraph().text = f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

s = prs.slides.add_slide(prs.slide_layouts[5])
s.shapes.title.text = "요약 KPI"
left = Inches(0.5); top = Inches(1.5); width = Inches(9); height = Inches(4)
rows = len(kpi)+1; cols = 2
tbl = s.shapes.add_table(rows, cols, left, top, width, height).table
tbl.cell(0,0).text = "지표"; tbl.cell(0,1).text = "값"
for r, (k,v) in enumerate(kpi.items(), start=1):
    tbl.cell(r,0).text = k
    if isinstance(v, (int, np.integer)): tbl.cell(r,1).text = f"{v:,}"
    elif isinstance(v, float): tbl.cell(r,1).text = f"{v:.2f}"
    else: tbl.cell(r,1).text = str(v)

def add_chart_slide(title, img):
    s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = title
    s.shapes.add_picture(img, Inches(0.7), Inches(1.2), height=Inches(5.2))

add_chart_slide("매매 vs 전세가 추이", IMG_PRICE_BOTH)
add_chart_slide("전세가율 추이", IMG_JEONSE_RATIO)
add_chart_slide("PIR 추이", IMG_PIR)
add_chart_slide("인구 추이", IMG_POP)
add_chart_slide("준공 물량(공급)", IMG_SUP)
add_chart_slide("일자리(취업자) 추이", IMG_JOBS)
add_chart_slide("경매 낙찰가율 추이", IMG_AUC)

news_slide = prs.slides.add_slide(prs.slide_layouts[5])
news_slide.shapes.title.text = "최근 1개월 뉴스 (청주시 하이닉스)"
headers = ["제목", "매체", "게재일", "링크"]
n = len(news_items)
ntbl = news_slide.shapes.add_table(rows=n+1, cols=4,
                                   left=Inches(0.3), top=Inches(1.3),
                                   width=Inches(9.4), height=Inches(4.8)).table
for c,h in enumerate(headers): ntbl.cell(0,c).text = h
for r,it in enumerate(news_items, start=1):
    ntbl.cell(r,0).text = it["title"]
    ntbl.cell(r,1).text = it["source"]
    ntbl.cell(r,2).text = it["pubdate"].strftime("%Y-%m-%d %H:%M")
    ntbl.cell(r,3).text = it["link"]

prs.save(PPT_PATH); conn.close(); print(PPT_PATH)
